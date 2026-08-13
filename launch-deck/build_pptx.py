#!/usr/bin/env python
"""Wrap the rendered slide PNGs into a 16:9 PowerPoint, one full-bleed image per slide."""
import glob, os
from pptx import Presentation
from pptx.util import Emu

here = os.path.dirname(os.path.abspath(__file__))
imgs = sorted(glob.glob(os.path.join(here, "shots", "slide-*.png")))
assert imgs, "no slide PNGs found"

# 16:9 widescreen
W, H = Emu(12192000), Emu(6858000)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank

for img in imgs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, 0, 0, width=W, height=H)

out = os.path.join(here, "allmeta-Ontology-v1.0-发布演讲.pptx")
prs.save(out)
print(f"PPTX written: {out}  ({len(imgs)} slides, {os.path.getsize(out)//1024} KB)")
